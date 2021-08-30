from pptx import Presentation
import os
from PIL import Image
#import water

def _add_image(slide, placeholder_id, image_url):
    placeholder = slide.placeholders[placeholder_id]
 
    im = Image.open(image_url)
    width, height = im.size
 
    placeholder.height = height
    placeholder.width = width
    placeholder = placeholder.insert_picture(image_url)
    image_ratio = width / height
    placeholder_ratio = placeholder.width / placeholder.height
    ratio_difference = placeholder_ratio - image_ratio

    if ratio_difference > 0:
        difference_on_each_side = ratio_difference / 2
        placeholder.crop_left = -difference_on_each_side
        placeholder.crop_right = -difference_on_each_side
    else:
        difference_on_each_side = -ratio_difference / 2
        placeholder.crop_bottom = -difference_on_each_side
        placeholder.crop_top = -difference_on_each_side

prs = Presentation()
layout8=prs.slide_layouts[8]

slide=prs.slides.add_slide(layout8)
print("It's make 5 slide only\n")
print("For Slide_1 ")

a = input("title\t:")
b = input("sub\t:")

c = input("image name\t:")
ppt_name = input("ppt name\t:")
title =  slide.shapes.title.text = a
sub = slide.placeholders[2].text = b
_add_image(slide,1,(c))

print("For Slide_2 \n")
layout8=prs.slide_layouts[8]

slide1=prs.slides.add_slide(layout8)
a1 = input("title\t:")
b1 = input("sub\t:")

c1 = input("image name\t:")
#ppt_name = input("ppt name\t")
title =  slide1.shapes.title.text = a1
sub = slide1.placeholders[2].text = b1
_add_image(slide1,1,(c1))


print("For Slide_3 \n")
layout8=prs.slide_layouts[8]

slide2=prs.slides.add_slide(layout8)
a2 = input("title\t:")
b2 = input("sub\t:")

c2 = input("image name\t:")
#ppt_name = input("ppt name\t")
title =  slide2.shapes.title.text = a2
sub = slide2.placeholders[2].text = b2
_add_image(slide2,1,(c2))


print("For Slide_4 \n")
layout8=prs.slide_layouts[8]

slide3=prs.slides.add_slide(layout8)
a3 = input("title\t:")
b3 = input("sub\t:")

c3 = input("image name\t:")
#ppt_name = input("ppt name\t")
title =  slide3.shapes.title.text = a3
sub = slide3.placeholders[2].text = b3
_add_image(slide3,1,(c3))



print("For Slide_5 \n")
layout8=prs.slide_layouts[8]

slide4=prs.slides.add_slide(layout8)
a4 = input("title\t:")
b4 = input("sub\t:")

c4 = input("image name\t:")
#ppt_name = input("ppt name\t")
title =  slide4.shapes.title.text = a4
sub = slide4.placeholders[2].text = b4
_add_image(slide4,1,(c4))


prs.save(ppt_name)

os.startfile(ppt_name)